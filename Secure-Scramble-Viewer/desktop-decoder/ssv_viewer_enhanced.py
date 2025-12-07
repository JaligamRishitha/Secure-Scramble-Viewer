#!/usr/bin/env python3
"""
SSV Enhanced Viewer - Full document viewing with zoom
Supports: Images (with zoom), PDF, Word, Excel, PowerPoint
"""

import sys
import os
import hashlib
import tkinter as tk
from tkinter import messagebox, scrolledtext, ttk
from pathlib import Path
from io import BytesIO
from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
from cryptography.hazmat.backends import default_backend
from cryptography.hazmat.primitives import padding

# Try to import optional libraries
try:
    from PIL import Image, ImageTk
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class SSVDecoder:
    """Standalone SSV file decoder"""
    
    def __init__(self, secret_key: str):
        self.key = hashlib.sha256(secret_key.encode()).digest()
    
    def decrypt_file(self, encrypted_data: bytes, salt: bytes, iv: bytes) -> bytes:
        """Decrypt file data using AES-256-CBC"""
        derived_key = hashlib.pbkdf2_hmac('sha256', self.key, salt, 100000, dklen=32)
        
        cipher = Cipher(
            algorithms.AES(derived_key),
            modes.CBC(iv),
            backend=default_backend()
        )
        decryptor = cipher.decryptor()
        padded_data = decryptor.update(encrypted_data) + decryptor.finalize()
        
        unpadder = padding.PKCS7(128).unpadder()
        original_data = unpadder.update(padded_data) + unpadder.finalize()
        
        return original_data
    
    def parse_ssv_file(self, ssv_data: bytes) -> tuple:
        """Parse .ssv file and return (decrypted_data, original_filename)"""
        version = ssv_data[0:4]
        salt = ssv_data[4:20]
        iv = ssv_data[20:36]
        fn_salt = ssv_data[36:52]
        fn_iv = ssv_data[52:68]
        filename_length = int.from_bytes(ssv_data[68:72], byteorder='big')
        
        encrypted_filename = ssv_data[72:72+filename_length]
        encrypted_data = ssv_data[72+filename_length:]
        
        original_filename = self.decrypt_file(encrypted_filename, fn_salt, fn_iv).decode('utf-8')
        original_data = self.decrypt_file(encrypted_data, salt, iv)
        
        return original_data, original_filename


class SSVViewerApp:
    """Enhanced viewer with zoom and full document support"""
    
    def __init__(self, root, ssv_file_path=None):
        self.root = root
        self.root.title("Secure Scramble Viewer")
        self.root.geometry("1200x800")
        
        self.secret_key = self.load_secret_key()
        self.decoder = SSVDecoder(self.secret_key)
        
        # Zoom state
        self.zoom_level = 1.0
        self.current_image = None
        self.current_data = None
        self.current_filename = None
        
        # PDF state
        self.pdf_document = None
        self.current_page = 0
        
        self.setup_ui()
        
        if ssv_file_path and os.path.exists(ssv_file_path):
            self.root.after(100, lambda: self.open_file(ssv_file_path))
    
    def load_secret_key(self) -> str:
        """Load secret key from config file"""
        config_file = Path.home() / ".ssv_decoder" / "config.txt"
        if config_file.exists():
            try:
                return config_file.read_text().strip()
            except:
                pass
        return "CHANGE-THIS-TO-A-SECURE-RANDOM-KEY-AT-LEAST-32-CHARACTERS-LONG"
    
    def setup_ui(self):
        """Setup the viewer UI"""
        # Header
        header = tk.Frame(self.root, bg="#1a1a2e", height=60)
        header.pack(fill=tk.X)
        header.pack_propagate(False)
        
        title = tk.Label(
            header,
            text="🔒 Secure Scramble Viewer",
            font=("Arial", 18, "bold"),
            bg="#1a1a2e",
            fg="#00d4ff"
        )
        title.pack(side=tk.LEFT, padx=20, pady=15)
        
        self.info_label = tk.Label(
            header,
            text="No file opened",
            font=("Arial", 10),
            bg="#1a1a2e",
            fg="#aaaaaa"
        )
        self.info_label.pack(side=tk.RIGHT, padx=20)
        
        # Toolbar
        self.toolbar = tk.Frame(self.root, bg="#0f3460", height=50)
        self.toolbar.pack(fill=tk.X)
        self.toolbar.pack_propagate(False)
        
        # Zoom controls (hidden by default)
        self.zoom_frame = tk.Frame(self.toolbar, bg="#0f3460")
        
        tk.Button(
            self.zoom_frame,
            text="🔍−",
            command=self.zoom_out,
            font=("Arial", 12, "bold"),
            bg="#00d4ff",
            fg="#1a1a2e",
            padx=10,
            pady=5,
            cursor="hand2",
            relief=tk.FLAT
        ).pack(side=tk.LEFT, padx=5)
        
        self.zoom_label = tk.Label(
            self.zoom_frame,
            text="100%",
            font=("Arial", 10),
            bg="#0f3460",
            fg="#ffffff"
        )
        self.zoom_label.pack(side=tk.LEFT, padx=10)
        
        tk.Button(
            self.zoom_frame,
            text="🔍+",
            command=self.zoom_in,
            font=("Arial", 12, "bold"),
            bg="#00d4ff",
            fg="#1a1a2e",
            padx=10,
            pady=5,
            cursor="hand2",
            relief=tk.FLAT
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Button(
            self.zoom_frame,
            text="Reset",
            command=self.zoom_reset,
            font=("Arial", 10),
            bg="#888888",
            fg="#ffffff",
            padx=10,
            pady=5,
            cursor="hand2",
            relief=tk.FLAT
        ).pack(side=tk.LEFT, padx=5)
        
        # PDF navigation (hidden by default)
        self.pdf_nav_frame = tk.Frame(self.toolbar, bg="#0f3460")
        
        tk.Button(
            self.pdf_nav_frame,
            text="◀ Previous",
            command=self.pdf_prev_page,
            font=("Arial", 10),
            bg="#00d4ff",
            fg="#1a1a2e",
            padx=10,
            pady=5,
            cursor="hand2",
            relief=tk.FLAT
        ).pack(side=tk.LEFT, padx=5)
        
        self.pdf_page_label = tk.Label(
            self.pdf_nav_frame,
            text="Page 1 / 1",
            font=("Arial", 10),
            bg="#0f3460",
            fg="#ffffff"
        )
        self.pdf_page_label.pack(side=tk.LEFT, padx=10)
        
        tk.Button(
            self.pdf_nav_frame,
            text="Next ▶",
            command=self.pdf_next_page,
            font=("Arial", 10),
            bg="#00d4ff",
            fg="#1a1a2e",
            padx=10,
            pady=5,
            cursor="hand2",
            relief=tk.FLAT
        ).pack(side=tk.LEFT, padx=5)
        
        # Content area with scrollbar
        self.content_container = tk.Frame(self.root, bg="#16213e")
        self.content_container.pack(fill=tk.BOTH, expand=True)
        
        self.canvas = tk.Canvas(self.content_container, bg="#16213e", highlightthickness=0)
        scrollbar = tk.Scrollbar(self.content_container, orient="vertical", command=self.canvas.yview)
        self.scrollable_frame = tk.Frame(self.canvas, bg="#16213e")
        
        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all"))
        )
        
        self.canvas_window = self.canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        self.canvas.configure(yscrollcommand=scrollbar.set)
        
        self.canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Make scrollable_frame expand to canvas width
        self.canvas.bind("<Configure>", self._on_canvas_configure)
        
        # Mouse wheel scrolling
        self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)
        
        self.show_welcome()
        self.root.bind("<Button-3>", lambda e: "break")
    
    def _on_mousewheel(self, event):
        self.canvas.yview_scroll(int(-1*(event.delta/120)), "units")
    
    def _on_canvas_configure(self, event):
        # Make scrollable_frame fill the canvas width so content centers properly
        self.canvas.itemconfig(self.canvas_window, width=event.width)
    
    def show_welcome(self):
        """Show welcome message"""
        self.hide_toolbars()
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        welcome = tk.Label(
            self.scrollable_frame,
            text="🔒 Secure Scramble Viewer\n\nDouble-click .ssv files to view them securely\n\n"
            "Supports: Images (with zoom), PDF, Word, Excel, PowerPoint\n\n"
            "🔒 View Only Mode - No Downloads",
            font=("Arial", 14),
            bg="#16213e",
            fg="#aaaaaa",
            justify=tk.CENTER
        )
        welcome.pack(expand=True, pady=200)
    
    def hide_toolbars(self):
        """Hide all toolbars"""
        self.zoom_frame.pack_forget()
        self.pdf_nav_frame.pack_forget()
    
    def open_file(self, file_path: str):
        """Open and display SSV file"""
        try:
            with open(file_path, 'rb') as f:
                ssv_data = f.read()
            
            original_data, original_filename = self.decoder.parse_ssv_file(ssv_data)
            
            self.current_data = original_data
            self.current_filename = original_filename
            
            file_size = len(original_data) / 1024
            self.info_label.config(
                text=f"📄 {original_filename} ({file_size:.1f} KB) - View Only"
            )
            
            self.display_content(original_data, original_filename)
            
        except Exception as e:
            messagebox.showerror(
                "Decryption Error",
                f"Failed to decrypt file:\n\n{str(e)}\n\nPossible causes:\n"
                f"1. Wrong secret key\n2. Corrupted file\n3. File encrypted with different key"
            )
            self.show_welcome()
    
    def display_content(self, data: bytes, filename: str):
        """Display content based on file type"""
        self.hide_toolbars()
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        ext = os.path.splitext(filename)[1].lower()
        
        if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
            self.display_image(data, filename)
        elif ext in ['.txt', '.md', '.log', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.c', '.cpp']:
            self.display_text(data, filename)
        elif ext == '.pdf':
            self.display_pdf(data, filename)
        elif ext in ['.docx', '.odt']:
            self.display_word(data, filename, ext)
        elif ext in ['.xlsx', '.ods']:
            self.display_excel(data, filename)
        elif ext in ['.pptx', '.odp']:
            self.display_powerpoint(data, filename)
        else:
            self.display_binary_info(data, filename)
    
    def display_image(self, data: bytes, filename: str):
        """Display image with zoom support"""
        if not PILLOW_AVAILABLE:
            self.show_library_missing("Pillow", "image viewing", "pip install Pillow")
            return
        
        try:
            self.current_image = Image.open(BytesIO(data))
            self.zoom_level = 1.0
            self.zoom_frame.pack(side=tk.LEFT, padx=20, pady=10)
            self.render_image()
        except Exception as e:
            self.show_error(f"Failed to display image: {str(e)}")
    
    def render_image(self):
        """Render image at current zoom level"""
        if not self.current_image:
            return
        
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        # Calculate size
        img_width, img_height = self.current_image.size
        new_width = int(img_width * self.zoom_level)
        new_height = int(img_height * self.zoom_level)
        
        # Resize
        resized = self.current_image.resize((new_width, new_height), Image.Resampling.LANCZOS)
        photo = ImageTk.PhotoImage(resized)
        
        # Display - centered
        label = tk.Label(self.scrollable_frame, image=photo, bg="#16213e")
        label.image = photo
        label.pack(expand=True, pady=20)
        
        # Update zoom label
        self.zoom_label.config(text=f"{int(self.zoom_level * 100)}%")
        
        # Watermark
        watermark = tk.Label(
            self.scrollable_frame,
            text=f"🔒 {self.current_filename} - View Only Mode",
            font=("Arial", 10),
            bg="#16213e",
            fg="#888888"
        )
        watermark.pack(pady=10)
    
    def zoom_in(self):
        """Zoom in"""
        if self.zoom_level < 3.0:
            self.zoom_level += 0.25
            self.render_image()
    
    def zoom_out(self):
        """Zoom out"""
        if self.zoom_level > 0.25:
            self.zoom_level -= 0.25
            self.render_image()
    
    def zoom_reset(self):
        """Reset zoom"""
        self.zoom_level = 1.0
        self.render_image()
    
    def display_pdf(self, data: bytes, filename: str):
        """Display PDF with page navigation"""
        if not PDF_AVAILABLE:
            self.show_library_missing("PyMuPDF", "PDF viewing", "pip install PyMuPDF")
            return
        
        try:
            self.pdf_document = fitz.open(stream=data, filetype="pdf")
            self.current_page = 0
            self.pdf_nav_frame.pack(side=tk.LEFT, padx=20, pady=10)
            self.render_pdf_page()
        except Exception as e:
            self.show_error(f"Failed to display PDF: {str(e)}")
    
    def render_pdf_page(self):
        """Render current PDF page"""
        if not self.pdf_document:
            return
        
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        page = self.pdf_document[self.current_page]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for clarity
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        photo = ImageTk.PhotoImage(img)
        label = tk.Label(self.scrollable_frame, image=photo, bg="#16213e")
        label.image = photo
        label.pack(expand=True, fill=tk.BOTH, pady=20)
        
        # Update page label
        total_pages = len(self.pdf_document)
        self.pdf_page_label.config(text=f"Page {self.current_page + 1} / {total_pages}")
        
        # Watermark
        watermark = tk.Label(
            self.scrollable_frame,
            text=f"🔒 {self.current_filename} - View Only Mode - Page {self.current_page + 1}/{total_pages}",
            font=("Arial", 10),
            bg="#16213e",
            fg="#888888"
        )
        watermark.pack(pady=10)
    
    def pdf_next_page(self):
        """Go to next PDF page"""
        if self.pdf_document and self.current_page < len(self.pdf_document) - 1:
            self.current_page += 1
            self.render_pdf_page()
    
    def pdf_prev_page(self):
        """Go to previous PDF page"""
        if self.pdf_document and self.current_page > 0:
            self.current_page -= 1
            self.render_pdf_page()
    
    def display_word(self, data: bytes, filename: str, ext: str = '.docx'):
        """Display Word document or ODT"""
        # Handle ODT files
        if ext == '.odt':
            self.display_odt(data, filename)
            return
            
        if not DOCX_AVAILABLE:
            self.show_library_missing("python-docx", "Word viewing", "pip install python-docx")
            return
        
        try:
            doc = Document(BytesIO(data))
            
            text_widget = scrolledtext.ScrolledText(
                self.scrollable_frame,
                wrap=tk.WORD,
                font=("Arial", 11),
                bg="#0f3460",
                fg="#ffffff",
                padx=15,
                pady=15
            )
            text_widget.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            
            for para in doc.paragraphs:
                text_widget.insert(tk.END, para.text + "\n\n")
            
            text_widget.config(state=tk.DISABLED)
            
            watermark = tk.Label(
                self.scrollable_frame,
                text=f"🔒 {filename} - View Only Mode",
                font=("Arial", 10),
                bg="#16213e",
                fg="#888888"
            )
            watermark.pack(pady=5)
            
        except Exception as e:
            # Try to display as text if docx parsing fails
            self.display_text(data, filename)
    
    def display_odt(self, data: bytes, filename: str):
        """Display ODT (OpenDocument Text) file"""
        try:
            import zipfile
            from xml.etree import ElementTree as ET
            
            with zipfile.ZipFile(BytesIO(data)) as zf:
                content = zf.read('content.xml')
                root = ET.fromstring(content)
                
                # Extract text from ODT
                text_content = []
                for elem in root.iter():
                    if elem.text:
                        text_content.append(elem.text)
                
                text_widget = scrolledtext.ScrolledText(
                    self.scrollable_frame,
                    wrap=tk.WORD,
                    font=("Arial", 11),
                    bg="#0f3460",
                    fg="#ffffff",
                    padx=15,
                    pady=15
                )
                text_widget.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
                text_widget.insert(tk.END, '\n'.join(text_content))
                text_widget.config(state=tk.DISABLED)
                
                watermark = tk.Label(
                    self.scrollable_frame,
                    text=f"🔒 {filename} - View Only Mode",
                    font=("Arial", 10),
                    bg="#16213e",
                    fg="#888888"
                )
                watermark.pack(pady=5)
        except Exception as e:
            self.show_error(f"Failed to display ODT: {str(e)}")
    
    def display_excel(self, data: bytes, filename: str):
        """Display Excel spreadsheet"""
        if not EXCEL_AVAILABLE:
            self.show_library_missing("openpyxl", "Excel viewing", "pip install openpyxl")
            return
        
        try:
            wb = openpyxl.load_workbook(BytesIO(data))
            
            # Create notebook for sheets
            notebook = ttk.Notebook(self.scrollable_frame)
            notebook.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                frame = tk.Frame(notebook, bg="#0f3460")
                text_widget = scrolledtext.ScrolledText(
                    frame,
                    wrap=tk.NONE,
                    font=("Consolas", 10),
                    bg="#0f3460",
                    fg="#ffffff",
                    padx=10,
                    pady=10
                )
                text_widget.pack(fill=tk.BOTH, expand=True)
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    text_widget.insert(tk.END, row_text + "\n")
                
                text_widget.config(state=tk.DISABLED)
                notebook.add(frame, text=sheet_name)
            
            watermark = tk.Label(
                self.scrollable_frame,
                text=f"🔒 {filename} - View Only Mode",
                font=("Arial", 10),
                bg="#16213e",
                fg="#888888"
            )
            watermark.pack(pady=10)
            
        except Exception as e:
            self.show_error(f"Failed to display Excel file: {str(e)}")
    
    def display_powerpoint(self, data: bytes, filename: str):
        """Display PowerPoint presentation"""
        if not PPTX_AVAILABLE:
            self.show_library_missing("python-pptx", "PowerPoint viewing", "pip install python-pptx")
            return
        
        try:
            prs = Presentation(BytesIO(data))
            
            text_widget = scrolledtext.ScrolledText(
                self.scrollable_frame,
                wrap=tk.WORD,
                font=("Arial", 11),
                bg="#0f3460",
                fg="#ffffff",
                padx=20,
                pady=20
            )
            text_widget.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
            
            for i, slide in enumerate(prs.slides, 1):
                text_widget.insert(tk.END, f"=== Slide {i} ===\n\n", "heading")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_widget.insert(tk.END, shape.text + "\n\n")
                text_widget.insert(tk.END, "\n" + "="*50 + "\n\n")
            
            text_widget.tag_config("heading", font=("Arial", 12, "bold"))
            text_widget.config(state=tk.DISABLED)
            
            watermark = tk.Label(
                self.scrollable_frame,
                text=f"🔒 {filename} - View Only Mode - {len(prs.slides)} slides",
                font=("Arial", 10),
                bg="#16213e",
                fg="#888888"
            )
            watermark.pack(pady=10)
            
        except Exception as e:
            self.show_error(f"Failed to display PowerPoint: {str(e)}")
    
    def display_text(self, data: bytes, filename: str):
        """Display text content"""
        try:
            try:
                text_content = data.decode('utf-8')
            except:
                text_content = data.decode('latin-1')
            
            text_widget = scrolledtext.ScrolledText(
                self.scrollable_frame,
                wrap=tk.WORD,
                font=("Consolas", 11),
                bg="#0f3460",
                fg="#ffffff",
                padx=15,
                pady=15
            )
            text_widget.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            text_widget.insert(1.0, text_content)
            text_widget.config(state=tk.DISABLED)
            
            watermark = tk.Label(
                self.scrollable_frame,
                text=f"🔒 {filename} - View Only Mode",
                font=("Arial", 10),
                bg="#16213e",
                fg="#888888"
            )
            watermark.pack(pady=5)
            
        except Exception as e:
            self.show_error(f"Failed to display text: {str(e)}")
    
    def display_binary_info(self, data: bytes, filename: str):
        """Show info for binary/unknown files"""
        ext = os.path.splitext(filename)[1].upper()
        
        info = tk.Label(
            self.scrollable_frame,
            text=f"📦 Binary File\n\n{filename}\n\nType: {ext or 'Unknown'}\n"
            f"Size: {len(data) / 1024:.1f} KB\n\n"
            f"This file type cannot be previewed.\n\n"
            f"🔒 View Only Mode - No Downloads",
            font=("Arial", 12),
            bg="#16213e",
            fg="#aaaaaa",
            justify=tk.CENTER
        )
        info.pack(expand=True, pady=200)
    
    def show_library_missing(self, library: str, feature: str, install_cmd: str):
        """Show message when library is missing"""
        info = tk.Label(
            self.scrollable_frame,
            text=f"📦 Library Required\n\n"
            f"{library} is required for {feature}\n\n"
            f"Install with:\n{install_cmd}\n\n"
            f"File: {self.current_filename}\n"
            f"Size: {len(self.current_data) / 1024:.1f} KB\n\n"
            f"The file is decrypted and secured.\n\n"
            f"🔒 View Only Mode - No Downloads",
            font=("Arial", 12),
            bg="#16213e",
            fg="#aaaaaa",
            justify=tk.CENTER
        )
        info.pack(expand=True, pady=200)
    
    def show_error(self, message: str):
        """Show error message"""
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        error = tk.Label(
            self.scrollable_frame,
            text=f"❌ Error\n\n{message}",
            font=("Arial", 12),
            bg="#16213e",
            fg="#e74c3c",
            justify=tk.CENTER
        )
        error.pack(expand=True, pady=200)


def main():
    """Main entry point"""
    ssv_file = None
    if len(sys.argv) > 1:
        ssv_file = sys.argv[1]
    
    root = tk.Tk()
    app = SSVViewerApp(root, ssv_file)
    root.protocol("WM_DELETE_WINDOW", root.destroy)
    root.mainloop()


if __name__ == "__main__":
    main()
